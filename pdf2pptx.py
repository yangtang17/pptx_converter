import argparse
import glob
import logging
import math
import os
import sys
import tempfile

import imageio
import pdf2image
import pptx
import pptx.util

TMP_DIR = './test_images/'


def main():
    args = configure_parser().parse_args()
    setup_logger(args.verbose)

    with tempfile.TemporaryDirectory() as images_directory:
        logging.debug('### Created temporary images directory: %s\n', images_directory)
        pdf_to_images(args.pdf_path, images_directory)
        images_to_pptx(images_directory, args.pptx_path)


def pdf_to_images(pdf_path, images_directory=None, dpi=600):
    logging.info('### Converting %s into images in %s, dpi=%s\n', pdf_path, images_directory, dpi)
    pil_images = pdf_to_pil(pdf_path, dpi)
    save_images(pil_images, images_directory)


def pdf_to_pil(pdf_path, dpi=600):
    # This method reads a pdf and converts it into a sequence of images
    # PDF_PATH sets the path to the PDF file
    # dpi parameter assists in adjusting the resolution of the image
    # output_folder parameter sets the path to the folder to which the PIL images can be stored (optional)
    # first_page parameter allows you to set a first page to be processed by pdftoppm
    # last_page parameter allows you to set a last page to be processed by pdftoppm
    # fmt parameter allows to set the format of pdftoppm conversion (PpmImageFile, TIFF)
    # thread_count parameter allows you to set how many thread will be used for conversion.
    # userpw parameter allows you to set a password to unlock the converted PDF
    # use_cropbox parameter allows you to use the crop box instead of the media box when converting
    # strict parameter allows you to catch pdftoppm syntax error with a custom type PDFSyntaxError

    pil_images = pdf2image.convert_from_path(pdf_path, dpi=dpi)
    return pil_images


def save_images(pil_images, save_dir):
    # This method helps in converting the images in PIL Image file format to the required image format
    index = 1
    digits = int(math.log10(len(pil_images))) + 1
    for image in pil_images:
        base_filename = 'page_' + str(index).zfill(digits) + '.jpg'
        file_path = os.path.join(save_dir, base_filename)
        logging.debug('Saving temporary image file: %s', file_path)

        image.save(file_path, 'JPEG')
        index += 1


def setup_logger(verbose):
    logging_level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=logging_level,
                        format='%(levelname)s - %(message)s',
                        handlers=[
                            logging.StreamHandler(sys.stdout)
                        ])


def images_to_pptx(images_directory, pptx_path):
    logging.info('### Converting images to %s, temporary image directory: %s', pptx_path, images_directory)
    prs = pptx.Presentation()
    for g in sorted(glob.glob(images_directory + '/*')):
        logging.debug('Converting %s', g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        img = imageio.imread(g)
        pic_height = prs.slide_height
        pic_width = min(prs.slide_width, int(pic_height * img.shape[1] / img.shape[0]))
        pic_left = (prs.slide_width - pic_width) / 2
        pic_top = 0
        slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

    prs.save(pptx_path)


def configure_parser():
    parser = argparse.ArgumentParser(prog='pdf2pptx', description='Script to convert pdf to pptx file')

    parser.add_argument('-f', '--from',
                        help='Pdf file path',
                        required=True,
                        dest='pdf_path')

    parser.add_argument('-t', '--to',
                        help='pptx file path',
                        required=True,
                        dest='pptx_path')

    parser.add_argument('-d', '--dpi',
                        help='DPI for the jpg images converted from PDF, defaulted to 600',
                        required=False,
                        dest='dpi')

    parser.add_argument("-v", "--verbose",
                        help="increase output verbosity",
                        action="store_true")
    return parser


if __name__ == '__main__':
    main()
